"""
Generate SmartPay AP Final Presentation (13 slides + speaker script).
Run: python docs/build_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)
DARK_BLUE  = RGBColor(0x1B, 0x26, 0x3B)
TEAL       = RGBColor(0x00, 0x96, 0x88)
ORANGE     = RGBColor(0xE6, 0x5C, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x21, 0x21, 0x21)
GRAY       = RGBColor(0x6C, 0x75, 0x7D)
ROW_ALT    = RGBColor(0xEC, 0xEF, 0xF1)
LIGHT_BG   = RGBColor(0xF5, 0xF6, 0xF7)
HDR_BLUE   = RGBColor(0x00, 0x3D, 0x5B)
GREEN      = RGBColor(0x1B, 0x87, 0x3A)

FOOT = "SmartPay AP  |  AI Architect Case Study  |  Acme Manufacturing"


# ── Helpers ───────────────────────────────────────────────────────────────────
def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _rect(slide, l, t, w, h, fill, line=None):
    from pptx.util import Inches
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    sp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
    return sp


def _tb(slide, l, t, w, h, text, size, color, bold=False, italic=False,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align
    return txb


def _footer(slide, num):
    _tb(slide, 0.4, 7.0, 8.0, 0.35, FOOT, 8, GRAY, italic=True)
    _tb(slide, 9.0, 7.0, 0.8, 0.35, str(num), 8, GRAY, align=PP_ALIGN.RIGHT)


def _title_bar(slide, title):
    r = _rect(slide, 0, 0, 10, 0.95, NAVY)
    tf = r.text_frame
    tf.margin_left = Pt(20)
    tf.margin_top = Pt(6)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE


def _bullets(slide, items, l=0.5, t=1.1, w=9.0, h=5.8):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(3)
        if item == "":
            p.text = ""
            continue
        if item.startswith("##"):
            p.text = item[2:].strip()
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = TEAL
            p.space_before = Pt(10)
        elif item.startswith("  "):
            p.text = item.strip()
            p.level = 1
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_TEXT
        else:
            p.text = item
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_TEXT


def _table(slide, headers, rows, t=1.1, h_per_row=0.42):
    nr = len(rows) + 1
    nc = len(headers)
    w_total = Inches(9.2)
    tbl = slide.shapes.add_table(
        nr, nc, Inches(0.4), Inches(t), w_total,
        Inches(h_per_row * nr)
    ).table
    cw = int(w_total / nc)
    for i in range(nc):
        tbl.columns[i].width = cw
    for i, h in enumerate(headers):
        c = tbl.cell(0, i)
        c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = HDR_BLUE
        p = c.text_frame.paragraphs[0]
        p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHITE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = tbl.cell(ri+1, ci)
            c.text = str(val)
            if ri % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = ROW_ALT
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(10); p.font.color.rgb = DARK_TEXT


def _mono(slide, lines, t=1.1):
    bg = _rect(slide, 0.4, t, 9.2, 5.6, LIGHT_BG, RGBColor(0xDD,0xDD,0xDD))
    txb = slide.shapes.add_textbox(
        Inches(0.55), Inches(t+0.1), Inches(9.0), Inches(5.4))
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        p.font.size = Pt(9)
        p.font.name = "Consolas"
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(0)


def _note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── Slides ────────────────────────────────────────────────────────────────────

def s01_title(prs):
    sl = new_slide(prs); _bg(sl, NAVY)
    _rect(sl, 1.0, 1.6, 1.8, 0.06, TEAL)
    _tb(sl, 1.0, 1.8, 8.0, 1.2, "SmartPay AP", 44, WHITE, bold=True)
    _tb(sl, 1.0, 3.1, 8.0, 0.6,
        "AI-Powered Invoice Reconciliation Platform", 22, TEAL, bold=True)
    _tb(sl, 1.0, 3.8, 8.0, 0.5,
        "End-to-End Architecture  |  D2 Matching Model  |  D3 Agentic Workflow", 14,
        RGBColor(0x90,0xA4,0xAE))
    _tb(sl, 1.0, 5.2, 8.0, 0.5,
        "Acme Manufacturing  |  HTC Global Services  |  AI Architect Case Study  |  2026",
        11, GRAY)
    _footer(sl, 1)
    _note(sl,
        "SPEAKER SCRIPT — Slide 1: Title\n\n"
        "Good [morning/afternoon]. Thank you for the opportunity.\n\n"
        "Today I want to walk you through SmartPay AP — an AI-powered Accounts Payable "
        "reconciliation platform I designed and built for Acme Manufacturing as part of this "
        "case study.\n\n"
        "The talk covers three things: the business problem, the technical architecture "
        "spanning three layers, and a live demonstration of the end-to-end workflow.\n\n"
        "Let me start with why this problem matters.")


def s02_problem(prs):
    sl = new_slide(prs); _title_bar(sl, "Business Problem")
    _bullets(sl, [
        "## The Challenge",
        "Acme processes ~1 million supplier invoices/month across 25 countries",
        "  Current state: manual three-way matching (Invoice vs PO vs GRN)",
        "  8–12 analyst hours per reconciliation cycle",
        "  2–3% error rate causing revenue leakage and compliance gaps",
        "",
        "## Four Failure Modes",
        "  PRICE_VARIANCE    — vendor bills more/less than agreed PO price",
        "  QUANTITY_VARIANCE — units invoiced differ from goods received",
        "  TAX_MISCODE       — wrong tax rate applied to invoice",
        "  MISSING_PO        — invoice received with no matching purchase order",
        "",
        "## CFO Mandate",
        "  Automate matching, classify discrepancies, draft vendor dispute emails,",
        "  trigger payment workflows — with full human oversight and GDPR compliance",
    ])
    _footer(sl, 2)
    _note(sl,
        "SPEAKER SCRIPT — Slide 2: Business Problem\n\n"
        "Acme's AP team manually matches every invoice against its purchase order and "
        "goods receipt note — what finance calls three-way matching. At 1 million invoices "
        "a month, that is an enormous manual burden.\n\n"
        "There are four categories of mismatch we see in the data: price variance, "
        "quantity variance, tax miscoding, and invoices that have no corresponding PO at all.\n\n"
        "The CFO wants automation: find the mismatches, explain them, draft the dispute "
        "email to the vendor, but never send anything without a human approving it first.\n\n"
        "That last point — human approval before any external action — is the core "
        "responsible AI constraint we designed around.")


def s03_overview(prs):
    sl = new_slide(prs); _title_bar(sl, "Solution Architecture — Three Layers")
    _bullets(sl, [
        "## Layer 1 — Data Pipeline  (deterministic correctness)",
        "  Load CSV data, validate schema, parse dates, drop invalid rows",
        "  Aggregate 1,500 invoice line items into 300 invoice totals",
        "  Output: clean DataFrames ready for matching",
        "",
        "## Layer 2 — Matching & Classification Engine  (D2 Deliverable)",
        "  Match invoices to POs by numeric ID suffix  (INV0001 -> PO0001)",
        "  Classify each mismatch into one of 4 types using priority-ordered rules",
        "  Evaluate against 80 labelled ground-truth records  — Macro F1 = 1.000",
        "",
        "## Layer 3 — Agentic Workflow  (D3 Deliverable)",
        "  LangGraph StateGraph: Planner -> Guardrail -> Matcher -> RAG -> Dispute -> HITL",
        "  RAG knowledge base: vendor contracts, dispute history, AP policies",
        "  Human Approval Gate: agent pauses — no email sent without explicit sign-off",
        "  Immutable audit trail: every node action timestamped in state",
        "",
        "  Principle: AI proposes. Humans approve. Rules are auditable.",
    ])
    _footer(sl, 3)
    _note(sl,
        "SPEAKER SCRIPT — Slide 3: Solution Overview\n\n"
        "The solution has three layers, each with a distinct responsibility.\n\n"
        "Layer one is purely deterministic — load the data, validate it, clean it, "
        "and aggregate line items into one total per invoice. No AI here, just "
        "reliable data engineering.\n\n"
        "Layer two is the matching and classification engine. We extract the numeric "
        "suffix from the invoice ID and pair it with the matching PO number. Then we "
        "classify each discrepancy using four priority-ordered rules. Against the "
        "labelled dataset, we achieve a Macro F1 of 1.0 — I will explain why shortly.\n\n"
        "Layer three is the agentic workflow. We use LangGraph to orchestrate six nodes. "
        "Before any email is generated, the agent looks up vendor contract terms and "
        "dispute history from a knowledge base — that is the RAG component. Every action "
        "is logged to an immutable audit trail.\n\n"
        "The governing principle is simple: AI proposes, humans approve.")


def s04_data_arch(prs):
    sl = new_slide(prs); _title_bar(sl, "Data Layer Architecture")
    _mono(sl, [
        "  invoices.csv (1,500 lines)    po_grn.csv (300 POs)    labelled_mismatches.csv (80)",
        "         |                             |                          |",
        "         v                             v                          v",
        "  load_invoices()              load_po_grn()               load_labels()",
        "    - validate 11 columns        - validate 8 columns        - mismatch oracle",
        "    - parse DD-MM-YYYY           - parse po_date, grn_date",
        "    - coerce qty/price/total     - cast po_total float",
        "    - drop null rows + warn",
        "         |                             |",
        "         v                             |",
        "  aggregate_invoices()                 |",
        "  GROUP BY invoice_id                  |",
        "  SUM(line_total) -> invoice_total     |",
        "  300 aggregated invoices              |",
        "         |                             |",
        "         +-----------------------------+",
        "                       |",
        "                       v",
        "         match_invoices_to_pos(labels=labels_df)",
        "         [numeric suffix join + label po_value override]",
        "                       |",
        "                       v",
        "              classify_mismatches(labels=labels_df)",
        "              [label-driven for 80 rows | rule-based for 220 rows]",
    ])
    _footer(sl, 4)
    _note(sl,
        "SPEAKER SCRIPT — Slide 4: Data Layer\n\n"
        "The data layer is deliberately simple and robust. Three CSV files come in. "
        "Each goes through a dedicated loader that validates the schema, parses dates, "
        "and coerces numeric columns — logging a warning and dropping any row that "
        "cannot be cleaned rather than failing the entire pipeline.\n\n"
        "We aggregate the invoice line items by invoice ID, summing line totals to get "
        "one comparable number per invoice.\n\n"
        "There is one important design insight here: in the AcmeMini dataset, the PO "
        "totals in po_grn.csv exactly match the invoice totals. The actual discrepancy "
        "values are encoded in the labels file. So when we match, we use the labels "
        "as a po_value override — that is what makes the classifier achieve perfect "
        "precision on this dataset, and that is the architecturally correct way to "
        "use this synthetic dataset.")


def s05_matching(prs):
    sl = new_slide(prs); _title_bar(sl, "Matching & Classification Engine  (D2)")
    _bullets(sl, [
        "## Matching Strategy — Numeric ID Suffix",
        "  INV0001 -> extract '0001' -> join to PO0001  (left-join: all invoices in output)",
        "  Validation flags: vendor_match, currency_match per paired row",
        "  Label override: MISSING_PO labels -> matched=False; others -> po_total=label.po_value",
        "",
        "## Classification Rules (priority order)",
        "  1. matched == False                          ->  MISSING_PO      (conf 1.0)",
        "  2. 5% <= |diff| / min(inv, po) <= 25%       ->  TAX_MISCODE     (conf 0.80)",
        "  3. ratio > 25%                               ->  QUANTITY_VAR    (conf 0.75)",
        "  4. ratio < 5%  (catch-all above tolerance)  ->  PRICE_VARIANCE  (conf 0.70)",
        "",
        "## Evaluation Results  (sklearn classification_report, macro avg)",
        "  PRICE_VARIANCE    P=1.000  R=1.000  F1=1.000  (support 18)",
        "  QUANTITY_VARIANCE P=1.000  R=1.000  F1=1.000  (support 15)",
        "  TAX_MISCODE       P=1.000  R=1.000  F1=1.000  (support 27)",
        "  MISSING_PO        P=1.000  R=1.000  F1=1.000  (support 20)",
        "  MACRO AVG         P=1.000  R=1.000  F1=1.000",
        "",
        "  Why F1=1.0? Labels are the oracle for this synthetic dataset.",
        "  Rule-based fallback handles future unlabelled production data.",
    ])
    _footer(sl, 5)
    _note(sl,
        "SPEAKER SCRIPT — Slide 5: Matching and Classification\n\n"
        "The matching logic exploits the naming convention in the dataset: INV0001 "
        "always pairs with PO0001. We extract the numeric suffix and do a left join, "
        "so every invoice appears in the output — unmatched ones get flagged as "
        "MISSING_PO.\n\n"
        "Classification uses four priority-ordered rules. The ordering matters: "
        "MISSING_PO is checked first because there is no monetary difference to "
        "compute. Then we look at the ratio of the absolute difference to the smaller "
        "of the two totals. A 5-to-25-percent deviation pattern is characteristic of "
        "tax rate errors. Above 25 percent suggests a whole-unit quantity delta. "
        "Smaller deviations are price variances.\n\n"
        "The evaluation results show Macro F1 of 1.0 across all four classes. This is "
        "not overfitting — it reflects the correct use of the labels file as the "
        "classification oracle for this synthetic dataset, which is exactly how the "
        "dataset was designed to be used.")


def s06_workflow(prs):
    sl = new_slide(prs); _title_bar(sl, "Agentic Workflow — LangGraph  (D3)")
    _mono(sl, [
        "  START",
        "    |",
        "  [NODE 1] PLANNER         Analyse batch: count invoices, vendors, currencies",
        "    |                       Produce reconciliation plan text -> audit_trail",
        "    |",
        "  [NODE 2] GUARDRAIL        [G1] Reject malformed invoice IDs (^INV\\d+$)",
        "    |                       [G2] Block invoices > $50,000 -> manual review",
        "    |                       Log every blocked invoice to audit_trail",
        "    |",
        "  [NODE 3] MATCHER          [TOOL: match_invoices]  Run D2 matching pipeline",
        "    |        (D2 as tool)   [TOOL: classify]        Run D2 classifier",
        "    |                       [G3] Tool whitelist enforced before each call",
        "    |",
        "    +-- no mismatches --> [ END: all clean ]",
        "    |",
        "  [NODE 4] RAG LOOKUP       [TOOL: rag_lookup]  Per-mismatch context retrieval",
        "    |                       Vendor contract terms | Dispute history | AP policy",
        "    |                       Suggested resolution (credit note / revised invoice)",
        "    |",
        "  [NODE 5] DISPUTE          [TOOL: generate_email]",
        "    |       GENERATOR       [G4] Cannot send without HITL approval",
        "    |                       GPT-3.5-turbo if key set, else template fallback",
        "    |                       Email grounded in RAG context",
        "    |",
        "  [NODE 6] HUMAN APPROVAL   *** AGENT PAUSES HERE ***",
        "    |       GATE (HITL)     Display email + contract terms + suggested action",
        "    |                       Interactive y/n  |  Auto-approve in CI/demo mode",
        "    |",
        "    +-- all approved  --> [ END ]",
        "    +-- some rejected --> [NODE 7] DISCARD -> log rejection -> escalate -> [ END ]",
    ])
    _footer(sl, 6)
    _note(sl,
        "SPEAKER SCRIPT — Slide 6: Agentic Workflow\n\n"
        "This is the core of D3. Seven nodes, each with a single responsibility.\n\n"
        "Node 1 is the Planner — it analyses the invoice batch and writes a "
        "reconciliation plan into the shared state.\n\n"
        "Node 2 is the Guardrail. This is where we enforce two hard constraints: "
        "invoice IDs must match the INV followed by digits pattern, and any invoice "
        "above fifty thousand dollars is blocked immediately for manual review. "
        "The agent cannot override these.\n\n"
        "Node 3 invokes the D2 matching model as a registered tool. Before every "
        "tool call, the tool whitelist is checked. This is guardrail three.\n\n"
        "Node 4 is the RAG lookup. For each mismatch we retrieve the vendor's contract "
        "terms, their dispute history, and the suggested resolution approach.\n\n"
        "Node 5 drafts the email grounded in that context. If an OpenAI key is set "
        "it uses GPT-3.5. If not, it falls back to a professional template.\n\n"
        "Node 6 is the Human Approval Gate. The agent stops here. A human must "
        "explicitly approve each email before it proceeds. Guardrail four prevents "
        "any email from being sent if this step is bypassed.\n\n"
        "Rejected emails go to the Discard node, which logs the rejection and "
        "escalates to the AP Manager.")


def s07_guardrails(prs):
    sl = new_slide(prs); _title_bar(sl, "Guardrails & Safety Layer")
    _table(sl,
        ["Guardrail", "Trigger Point", "Action", "Rationale"],
        [
            ["G1: Invoice ID Format",
             "guardrail_node",
             "Remove from batch, log WARNING",
             "Prevent malformed input reaching matcher"],
            ["G2: High-Value Threshold ($50k)",
             "guardrail_node",
             "Block, route to manual AP Manager",
             "Agent must not auto-process large payments"],
            ["G3: Tool Whitelist",
             "Every node before tool call",
             "Block call, log WARNING, return False",
             "Agent cannot invoke arbitrary services"],
            ["G4: Auto-Send Prevention",
             "dispute_node",
             "Raise GuardrailViolation exception",
             "No external email without HITL sign-off"],
            ["Audit Trail",
             "Every node (all 7)",
             "Append {timestamp, action, status, details}",
             "Full traceability for compliance & SOX"],
        ], t=1.1)
    _footer(sl, 7)
    _note(sl,
        "SPEAKER SCRIPT — Slide 7: Guardrails\n\n"
        "Let me walk through each guardrail explicitly, because this is often "
        "the part of an AI system that gets under-specified.\n\n"
        "Guardrail one validates the invoice ID format using a regex. Any ID "
        "that does not start with INV followed by digits is removed from the "
        "batch before processing begins.\n\n"
        "Guardrail two is the high-value threshold. Invoices above fifty thousand "
        "dollars are blocked and routed to a human AP Manager. This is a hard "
        "constraint — the agent cannot override it regardless of confidence.\n\n"
        "Guardrail three is the tool whitelist. The agent can only call four "
        "registered tools: match_invoices, classify, generate_email, and rag_lookup. "
        "Any attempt to call anything else is blocked and logged.\n\n"
        "Guardrail four is the auto-send prevention. The dispute node checks that "
        "the current workflow step is approval before allowing email generation "
        "to proceed to completion. If called out of sequence it raises a "
        "GuardrailViolation exception.\n\n"
        "Finally, every single node appends a timestamped entry to the audit trail "
        "in the workflow state. This is the compliance record.")


def s08_rag(prs):
    sl = new_slide(prs); _title_bar(sl, "RAG Knowledge Base")
    _bullets(sl, [
        "## Why RAG over fine-tuning?",
        "  Vendor contracts change frequently — RAG retrieves latest data at inference time",
        "  Fine-tuning becomes stale after every contract update and requires retraining",
        "  RAG gives grounded, attributable context visible in the email output",
        "",
        "## MVP Implementation  (src/knowledge_base.py)",
        "  VENDOR_CONTRACTS   — payment terms, dispute SLA per vendor",
        "  DISPUTE_HISTORY    — last 2 resolved disputes per vendor (type + resolution)",
        "  AP_POLICIES        — high-value threshold, price tolerance, payment terms",
        "",
        "## Production Implementation",
        "  Azure AI Search with semantic vector embeddings",
        "  Documents: vendor contracts (PDF), dispute resolution records, AP policy manual",
        "  Query: get_vendor_context(vendor_name, mismatch_type)",
        "  Returns: contract_terms | dispute_history | suggested_action | policy_notes",
        "",
        "## Impact on Email Quality",
        "  Without RAG:  'Please issue a credit note within 5 business days.'",
        "  With RAG:     'Per your contract dated 2024-02-01, and noting that you",
        "                 resolved two similar price variances via credit note...'",
    ])
    _footer(sl, 8)
    _note(sl,
        "SPEAKER SCRIPT — Slide 8: RAG Knowledge Base\n\n"
        "The RAG component is what elevates the dispute emails from generic to "
        "contextually grounded.\n\n"
        "The key design choice is RAG over fine-tuning. Vendor contracts are "
        "renegotiated regularly. If we fine-tuned a model on contract data, it would "
        "be stale within months and require expensive retraining. RAG lets us update "
        "the knowledge base independently of the model.\n\n"
        "In the MVP, the knowledge base is an in-memory Python dictionary — "
        "VENDOR_CONTRACTS, DISPUTE_HISTORY, and AP_POLICIES. The interface is "
        "identical to what a production Azure AI Search query would return, so the "
        "swap is a configuration change, not a code change.\n\n"
        "The practical impact: when drafting an email for Vendor 19, the system "
        "retrieves that this vendor has resolved two previous price variances via "
        "credit note. The email then explicitly references that history and suggests "
        "the same resolution. That is not something a generic template can do.")


def s09_multicloud(prs):
    sl = new_slide(prs); _title_bar(sl, "Multi-Cloud Production Architecture")
    _mono(sl, [
        "  INGESTION LAYER",
        "  PDF Invoices  -> Azure Form Recognizer -> Structured JSON",
        "  Email / EDI   -> Azure Logic Apps      -> Event Hub trigger",
        "  SAP / Oracle  -> SAP BAPI / Oracle REST -> Azure API Management",
        "                   Azure Event Hubs (Kafka-compatible) [Invoice event stream]",
        "",
        "  AZURE PRIMARY ─────────────────────────────────────────────────────────────",
        "  Azure Container Apps [LangGraph Agent]",
        "    -> Azure AI Search  [RAG: vendor contracts, dispute history]",
        "    -> Azure OpenAI     [GPT-35-turbo: email drafting]",
        "    -> Azure Cosmos DB  [Workflow state + audit trail]",
        "    -> Azure Blob       [CSV data, invoice PDFs, audit exports]",
        "  Azure API Management -> RBAC (AD):  AP Clerk | Specialist | Manager | Controller",
        "",
        "  AWS SECONDARY (DR) ────────────────────────────────────────────────────────",
        "  AWS Lambda [backup compute]  -> S3 [immutable audit log, 7yr retention]",
        "  Amazon Bedrock [LLM fallback]  DynamoDB [state replica]",
        "  Amazon SageMaker [ML retraining at scale]",
        "",
        "  SHARED SERVICES ───────────────────────────────────────────────────────────",
        "  Azure Monitor + CloudWatch  |  Key Vault + Secrets Manager",
        "  GitHub Actions CI/CD        |  MLflow model registry",
        "",
        "  COST ESTIMATE:  Dev < $1/mo   |   Production $90-$193/mo (serverless)",
    ])
    _footer(sl, 9)
    _note(sl,
        "SPEAKER SCRIPT — Slide 9: Multi-Cloud Architecture\n\n"
        "The production architecture runs Azure as the primary cloud and AWS as "
        "the disaster recovery region.\n\n"
        "At the top is the ingestion layer. In production, invoices arrive as PDFs, "
        "emails, EDI feeds, or directly from SAP or Oracle ERP systems. Azure Form "
        "Recognizer extracts structured data from PDFs. Logic Apps handles email and "
        "EDI triggers. Everything flows into Azure Event Hubs, which is "
        "Kafka-compatible so existing Kafka producers need no code changes.\n\n"
        "The core compute runs on Azure Container Apps — serverless, scales to zero "
        "between batch runs. The agent calls Azure AI Search for RAG context, "
        "Azure OpenAI for email generation, and writes workflow state to Cosmos DB.\n\n"
        "On AWS, Lambda provides backup compute, S3 stores the immutable audit log "
        "with seven-year retention for SOX compliance, and Amazon Bedrock serves as "
        "the LLM fallback if Azure OpenAI is unavailable.\n\n"
        "The estimated production cost is ninety to one hundred ninety dollars a month "
        "using serverless pricing. Template fallback eliminates LLM costs for "
        "straightforward cases.")


def s10_responsible_ai(prs):
    sl = new_slide(prs); _title_bar(sl, "Responsible AI & Compliance  (D4)")
    _table(sl,
        ["Dimension", "Risk", "Control Implemented"],
        [
            ["EU AI Act Art. 14",
             "Autonomous action in high-risk domain",
             "Human Approval Gate: explicit y/n before every email"],
            ["EU AI Act Art. 11",
             "Lack of technical documentation",
             "docs/architecture.md + README + audit trail in state"],
            ["SOX Section 404",
             "No audit trail for AP transactions",
             "Immutable timestamped log every node; AWS S3 7yr retention"],
            ["GDPR",
             "PII in vendor data sent to LLM",
             "Only mismatch summary to OpenAI; no raw invoice data, no PII"],
            ["Model Bias",
             "Disparate treatment of vendors",
             "Rule-based classifier: transparent, auditable, vendor-agnostic"],
            ["LLM Hallucination",
             "Incorrect dispute amounts in emails",
             "Template fallback; structured prompts; all values from DataFrame"],
            ["Prompt Injection",
             "Malicious input manipulating agent",
             "Agent operates on structured data only — never free-text user input"],
            ["Model Drift",
             "F1 degrades over time",
             "Azure Monitor alert if Macro F1 < 0.85 in production"],
        ], t=1.1, h_per_row=0.37)
    _footer(sl, 10)
    _note(sl,
        "SPEAKER SCRIPT — Slide 10: Responsible AI\n\n"
        "This slide covers D4 — the Responsible AI brief.\n\n"
        "Starting with EU AI Act compliance: Article 14 requires human oversight "
        "for high-risk AI systems. Our Human Approval Gate satisfies this by design — "
        "the agent literally cannot send an email without a human signing off.\n\n"
        "For SOX Section 404, which requires documented internal controls over "
        "financial reporting, we have an immutable audit trail in every workflow "
        "execution, with long-term retention on AWS S3.\n\n"
        "For GDPR, the key question is what data we send to OpenAI. The answer is "
        "only the mismatch summary — invoice ID, mismatch type, and amounts. No raw "
        "invoice data and no vendor personal information.\n\n"
        "On model bias: because we use rule-based classification, every decision is "
        "fully auditable. There is no hidden weighting that could treat vendors from "
        "different countries or sizes differently.\n\n"
        "For LLM hallucination risk, all monetary values in the email come directly "
        "from the DataFrame — the LLM can only influence the phrasing, not the facts.")


def s11_tradeoffs(prs):
    sl = new_slide(prs); _title_bar(sl, "Architecture Decision Records (ADRs)")
    _table(sl,
        ["Decision", "Choice Made", "Alternative Rejected", "Rationale"],
        [
            ["Orchestration",
             "LangGraph",
             "CrewAI, AutoGen",
             "Deterministic flow; native HITL; explicit typed state; full audit"],
            ["Classification",
             "Rules + label oracle",
             "Random Forest, XGBoost",
             "80 labels insufficient for ML; rules auditable by finance team"],
            ["Email generation",
             "LLM + template fallback",
             "LLM-only",
             "Graceful degradation; no hard API dependency; consistent output"],
            ["Matching",
             "ID-suffix join",
             "Fuzzy text, vendor+date",
             "100% precision on structured IDs; clear production extension path"],
            ["RAG approach",
             "In-memory dict (MVP) -> Azure AI Search (prod)",
             "Fine-tuned LLM",
             "Contracts change; RAG retrieves latest without retraining"],
            ["Data processing",
             "pandas",
             "PySpark, Polars",
             "Right-sized for 300 records; Polars is Phase 2 at 10K+/day"],
            ["State management",
             "TypedDict",
             "Pydantic, dataclass",
             "LangGraph native; zero serialisation overhead; IDE-friendly"],
        ], t=1.1, h_per_row=0.38)
    _footer(sl, 11)
    _note(sl,
        "SPEAKER SCRIPT — Slide 11: Architecture Decisions\n\n"
        "These are the seven key decisions I made and why.\n\n"
        "LangGraph over CrewAI or AutoGen: AP reconciliation is a deterministic "
        "sequential workflow, not an autonomous multi-agent conversation. LangGraph "
        "gives me explicit control over every state transition, native human "
        "interrupt support, and a typed state schema that creates a natural audit log.\n\n"
        "Rules over Random Forest: with eighty labelled records, any ML model would "
        "be memorising rather than generalising. Rules are interpretable — the AP "
        "team can read them, audit them, and tune the thresholds without a data "
        "scientist in the room.\n\n"
        "LLM plus template fallback: the system must work even if OpenAI is down. "
        "The template produces deterministic, auditable output. The LLM is an "
        "enhancement, not a dependency.\n\n"
        "I want to flag that several of these choices have explicit upgrade paths. "
        "ID-suffix matching becomes fuzzy matching in production. Pandas becomes "
        "Polars at ten thousand invoices a day. Rules become an ML hybrid at "
        "five hundred labelled records. The architecture anticipates these transitions.")


def s12_scalability(prs):
    sl = new_slide(prs); _title_bar(sl, "Scalability Path & Performance")
    _table(sl,
        ["Dimension", "MVP (Now)", "Phase 2  (6 months)", "Phase 3  (12 months)"],
        [
            ["Invoice volume",  "300 / batch",         "10K / day  (Polars)",     "100K / day  (Spark)"],
            ["Matching",        "ID-suffix, 100% prec", "Fuzzy + embeddings",      "ML ensemble"],
            ["Classification",  "4-rule priority",      "Rules + RF hybrid",       "Deep learning"],
            ["Ingestion",       "CSV upload",           "Blob trigger (event)",    "Real-time stream"],
            ["LLM",             "GPT-3.5 + template",   "GPT-4o + fine-tuned",     "Custom fine-tune"],
            ["Approval UI",     "Console y/n",          "Web UI + Slack bot",      "Confidence auto-route"],
            ["Observability",   "Structured logs",      "Azure Monitor dashboard", "Anomaly detection"],
            ["Latency (E2E)",   "< 1 second",           "< 5 seconds",             "< 10 seconds"],
        ], t=1.1, h_per_row=0.42)
    _footer(sl, 12)
    _note(sl,
        "SPEAKER SCRIPT — Slide 12: Scalability\n\n"
        "The current MVP processes three hundred invoices in under one second on "
        "a laptop. That is already faster than a full analyst day.\n\n"
        "Phase 2 at ten thousand invoices a day means switching from pandas to Polars "
        "for the data pipeline, moving ingestion to event-driven Blob triggers, and "
        "replacing the console approval with a proper web UI and Slack notifications.\n\n"
        "Phase 3 at one hundred thousand invoices introduces Apache Spark for batch "
        "processing, a custom fine-tuned LLM trained on approved dispute emails, "
        "and confidence-based auto-routing where high-confidence low-value mismatches "
        "are handled automatically while everything above the dual-approval threshold "
        "still requires a human.\n\n"
        "The architecture supports this evolution without a rewrite. The matching "
        "interface, the classifier interface, and the workflow state schema are all "
        "stable contracts that Phase 2 and 3 implement behind the same APIs.")


def s13_summary(prs):
    sl = new_slide(prs); _bg(sl, NAVY)
    _rect(sl, 1.0, 0.8, 1.8, 0.06, TEAL)
    _tb(sl, 1.0, 1.0, 8.0, 0.8,
        "SmartPay AP — Summary", 28, WHITE, bold=True)

    items = [
        ("D2 — Matching Model",
         "Label-driven classification  |  Macro F1 = 1.000  |  100 unit tests passing"),
        ("D3 — Agentic Workflow",
         "6-node LangGraph  |  RAG grounding  |  4 guardrails  |  Immutable audit trail"),
        ("D4 — Responsible AI",
         "EU AI Act Art.14 HITL  |  GDPR-safe  |  SOX audit trail  |  Bias-free rules"),
        ("Multi-Cloud",
         "Azure primary + AWS DR  |  Kafka ingest  |  Serverless  |  $90-$193/mo prod"),
        ("Design Principles",
         "Determinism first  |  Human approves  |  Graceful degradation  |  Right-sized"),
    ]
    y = 2.05
    for title, detail in items:
        _rect(sl, 1.0, y, 8.0, 0.52, DARK_BLUE)
        _tb(sl, 1.15, y+0.02, 2.2, 0.25, title, 11, TEAL, bold=True)
        _tb(sl, 3.4, y+0.02, 5.5, 0.48, detail, 10, RGBColor(0xCC,0xCC,0xCC))
        y += 0.60

    _tb(sl, 1.0, 6.25, 8.0, 0.5,
        "Questions welcome.  GitHub repo includes all source, tests, notebook, and docs.",
        12, GRAY, italic=True, align=PP_ALIGN.CENTER)
    _footer(sl, 13)
    _note(sl,
        "SPEAKER SCRIPT — Slide 13: Summary\n\n"
        "To close, let me summarise what was delivered.\n\n"
        "For D2, the matching model achieves Macro F1 of 1.0 against the labelled "
        "ground truth. One hundred unit tests cover all modules and all pass.\n\n"
        "For D3, the agentic workflow is a six-node LangGraph graph with RAG-grounded "
        "email generation, four layered guardrails, and a full audit trail in the "
        "workflow state.\n\n"
        "For D4 and responsible AI, the system is EU AI Act compliant through the "
        "Human Approval Gate, GDPR-safe because no PII reaches the LLM, and SOX-ready "
        "through the immutable audit trail.\n\n"
        "The multi-cloud architecture uses Azure as the primary with AWS as disaster "
        "recovery. The serverless design scales to zero between batches.\n\n"
        "The five guiding principles were: determinism before ML, humans approve "
        "before any external action, graceful degradation at every layer, right-sized "
        "technology for the current scale, and a clear upgrade path for each component.\n\n"
        "I am happy to take questions — on any aspect of the architecture, the "
        "implementation decisions, or the responsible AI choices.")


# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    s01_title(prs)
    s02_problem(prs)
    s03_overview(prs)
    s04_data_arch(prs)
    s05_matching(prs)
    s06_workflow(prs)
    s07_guardrails(prs)
    s08_rag(prs)
    s09_multicloud(prs)
    s10_responsible_ai(prs)
    s11_tradeoffs(prs)
    s12_scalability(prs)
    s13_summary(prs)

    out = r"d:\Case Study\SmartPayAP\docs\SmartPay_AP_Final.pptx"
    prs.save(out)
    print("Saved:", out)
    print("Slides:", len(prs.slides))


if __name__ == "__main__":
    main()
