"""
Generate Advanced Architecture Presentation for SmartPay AP.

Creates a professional 15-slide PPTX with:
- C4 model architecture views
- Industry-aligned three-way matching patterns
- LangGraph agentic workflow deep-dive
- EU AI Act compliance mapping
- HITL governance framework
- Production deployment topology
- Observability & SRE considerations
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette ────────────────────────────────────────────────────
NAVY = RGBColor(0x0D, 0x1B, 0x2A)
DARK_BLUE = RGBColor(0x1B, 0x26, 0x3B)
ACCENT_TEAL = RGBColor(0x00, 0x96, 0x88)
ACCENT_ORANGE = RGBColor(0xE6, 0x5C, 0x00)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x21, 0x21, 0x21)
MID_GRAY = RGBColor(0x6C, 0x75, 0x7D)
LIGHT_TEAL = RGBColor(0xE0, 0xF2, 0xF1)
TABLE_HEADER = RGBColor(0x00, 0x3D, 0x5B)
TABLE_ALT = RGBColor(0xEC, 0xEF, 0xF1)


def _add_footer(slide, text="SmartPay AP | AI Architect Case Study | Acme Manufacturing"):
    """Add a subtle footer to a slide."""
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(6.9), Inches(7), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = MID_GRAY
    p.font.italic = True


def _add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(9.2), Inches(6.9), Inches(0.6), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(9)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.8), Inches(1.5), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_TEAL
    line.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SmartPay AP"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "AI-Powered Invoice Reconciliation Architecture"
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT_TEAL

    p3 = tf2.add_paragraph()
    p3.text = "End-to-End System Design & Implementation"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0x90, 0xA4, 0xAE)
    p3.space_before = Pt(8)

    # Meta
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.4), Inches(8), Inches(1))
    tf3 = txBox3.text_frame
    p4 = tf3.paragraphs[0]
    p4.text = "Acme Manufacturing  |  AI Architect Assessment  |  June 2026"
    p4.font.size = Pt(12)
    p4.font.color.rgb = MID_GRAY

    _add_slide_number(slide, 1)
    return slide


def add_section_slide(prs, section_title, section_subtitle, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BLUE

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.0), Inches(2), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_TEAL
    line.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = section_subtitle
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0xB0, 0xBE, 0xC5)

    _add_slide_number(slide, num)
    return slide


def add_content_slide(prs, title, bullets, num, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_top = Pt(10)
    tf.margin_left = Pt(24)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content area
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(8.8), Inches(5.4))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]

        # Handle hierarchy
        if bullet.startswith("    "):
            p.text = bullet.strip()
            p.level = 2
            p.font.size = Pt(12)
            p.font.color.rgb = MID_GRAY
        elif bullet.startswith("  "):
            p.text = bullet.strip()
            p.level = 1
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_TEXT
        elif bullet.startswith("##"):
            p.text = bullet.replace("##", "").strip()
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = ACCENT_TEAL
            p.space_before = Pt(12)
        elif bullet == "":
            p.text = ""
            p.space_after = Pt(4)
        else:
            p.text = bullet
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(4)

    _add_footer(slide)
    _add_slide_number(slide, num)
    return slide


def add_table_slide(prs, title, headers, rows, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_top = Pt(10)
    tf.margin_left = Pt(24)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    left = Inches(0.4)
    top = Inches(1.3)
    width = Inches(9.2)
    row_height = Inches(0.45)
    height = row_height * num_rows

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Column widths
    col_width = int(width / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_width

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # Data
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_TEXT
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT

    _add_footer(slide)
    _add_slide_number(slide, num)
    return slide


def add_diagram_slide(prs, title, diagram_lines, num, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_top = Pt(10)
    tf.margin_left = Pt(24)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Diagram background box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.2), Inches(9.2), Inches(5.2)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    box.line.color.rgb = RGBColor(0xDE, 0xDE, 0xDE)

    # Diagram text
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.8), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(diagram_lines):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(10)
        p.font.name = "Consolas"
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(1)

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(9), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(10)
        p2.font.italic = True
        p2.font.color.rgb = MID_GRAY

    _add_footer(slide)
    _add_slide_number(slide, num)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ═══════════════════════════════════════════════════════════════════
    add_title_slide(prs)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 2: Executive Summary & Business Context
    # ═══════════════════════════════════════════════════════════════════
    add_content_slide(prs, "Executive Summary & Business Context", [
        "## Problem Domain",
        "Accounts Payable reconciliation at Acme Manufacturing: 300+ invoices/cycle,",
        "manual three-way matching (Invoice ↔ PO ↔ GRN) consuming 8–12 analyst hours",
        "",
        "## Solution",
        "SmartPay AP: An AI-powered reconciliation system combining deterministic",
        "matching rules with agentic orchestration and human-in-the-loop governance",
        "",
        "## Key Outcomes",
        "  • <1 second end-to-end matching (vs. 8–12 hours manual)",
        "  • 4-class mismatch taxonomy: PRICE, QUANTITY, TAX, MISSING_PO",
        "  • Zero autonomous external actions (HITL gate on all emails)",
        "  • EU AI Act Article 14 compliant (human oversight by design)",
        "  • Production-ready: guardrails, audit trail, graceful degradation",
    ], 2)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 3: C4 Model — System Context (Level 1)
    # ═══════════════════════════════════════════════════════════════════
    add_diagram_slide(prs, "C4 Model — Level 1: System Context", [
        "┌─────────────────────────────────────────────────────────────────────────┐",
        "│                         SYSTEM CONTEXT                                   │",
        "├─────────────────────────────────────────────────────────────────────────┤",
        "│                                                                          │",
        "│   ┌──────────────┐          ┌──────────────────┐          ┌──────────┐  │",
        "│   │   Finance    │          │   SmartPay AP    │          │  OpenAI  │  │",
        "│   │   Analyst    │─────────▶│   System         │─────────▶│  API     │  │",
        "│   │   [Person]   │◀─────────│   [Software]     │◀─────────│  [Ext]   │  │",
        "│   └──────────────┘  Reviews └──────────────────┘  LLM     └──────────┘  │",
        "│        │            & approves       │           emails                   │",
        "│        │             emails          │                                    │",
        "│        │                             ▼                                    │",
        "│        │                    ┌──────────────────┐                          │",
        "│        │                    │   ERP / Data     │                          │",
        "│        └───────────────────▶│   Source (CSV)   │                          │",
        "│           Uploads data      │   [External]     │                          │",
        "│                             └──────────────────┘                          │",
        "│                                                                          │",
        "└─────────────────────────────────────────────────────────────────────────┘",
    ], 3, "C4 Level 1: Who uses the system and what external systems does it interact with?")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 4: C4 Model — Container (Level 2)
    # ═══════════════════════════════════════════════════════════════════
    add_diagram_slide(prs, "C4 Model — Level 2: Container Diagram", [
        "┌───────────────────────── SmartPay AP ──────────────────────────────────┐",
        "│                                                                         │",
        "│  ┌────────────────┐   ┌────────────────┐   ┌────────────────────────┐  │",
        "│  │  Data Pipeline │   │  Matching       │   │  Agentic Workflow      │  │",
        "│  │  [Container]   │──▶│  Engine         │──▶│  [Container]           │  │",
        "│  │                │   │  [Container]    │   │                        │  │",
        "│  │  • data_loader │   │  • matcher      │   │  • LangGraph StateGraph│  │",
        "│  │  • CSV parsing │   │  • classifier   │   │  • Planner node        │  │",
        "│  │  • aggregation │   │  • evaluator    │   │  • Matcher tool node   │  │",
        "│  └────────────────┘   └────────────────┘   │  • Dispute node        │  │",
        "│                                             │  • Approval gate (HITL)│  │",
        "│  ┌────────────────┐   ┌────────────────┐   └────────────────────────┘  │",
        "│  │  Email Engine  │   │  Guardrails    │                                │",
        "│  │  [Container]   │   │  [Container]   │   ┌────────────────────────┐  │",
        "│  │                │   │                 │   │  Evaluation Engine     │  │",
        "│  │  • LLM mode    │   │  • ID validator │   │  [Container]           │  │",
        "│  │  • Template    │   │  • Tool whitelist│  │  • sklearn metrics    │  │",
        "│  │    fallback    │   │  • HITL enforcer│   │  • P/R/F1 per class   │  │",
        "│  └────────────────┘   └────────────────┘   └────────────────────────┘  │",
        "│                                                                         │",
        "└─────────────────────────────────────────────────────────────────────────┘",
    ], 4, "C4 Level 2: The major containers (deployable units) within SmartPay AP")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 5: Three-Layer Architecture Deep-Dive
    # ═══════════════════════════════════════════════════════════════════
    add_content_slide(prs, "Three-Layer Architecture (Industry Pattern)", [
        "## Layer 1: Deterministic Pipeline (Data Correctness)",
        "  • Schema validation, date parsing, numeric coercion",
        "  • Aggregation: 1,500 line items → 300 invoice totals",
        "  • ID-suffix matching (INV0001 → PO0001) — 100% precision",
        "",
        "## Layer 2: Inference Layer (Classification & Interpretation)",
        "  • Priority-ordered rule engine (4 mismatch classes)",
        "  • Confidence scoring per classification (0.70 – 1.0)",
        "  • Configurable thresholds (tolerance, quantity ratio)",
        "",
        "## Layer 3: State & Orchestration Layer (Agentic Control)",
        "  • LangGraph StateGraph with typed state schema",
        "  • Durable state transitions: planner → matcher → dispute → approval",
        "  • Conditional edges: approved → END, rejected → discard + audit log",
        "",
        "  This three-layer pattern aligns with production reconciliation architectures",
        "  (ref: Sphere Inc. Payment Reconciliation Agent Architecture, 2025)",
    ], 5)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 6: Three-Way Matching — Domain Model
    # ═══════════════════════════════════════════════════════════════════
    add_diagram_slide(prs, "Three-Way Matching — Domain Model", [
        "    ┌──────────────────────────────────────────────────────────────┐",
        "    │              THREE-WAY MATCHING CONTROL                       │",
        "    │                                                              │",
        "    │    ┌─────────────┐                                           │",
        "    │    │  PURCHASE   │  ← What was ordered?                      │",
        "    │    │  ORDER (PO) │    (items, qty, price, vendor, currency)  │",
        "    │    └──────┬──────┘                                           │",
        "    │           │                                                  │",
        "    │           │ Match by numeric ID suffix                       │",
        "    │           │ (INV000X ↔ PO000X)                               │",
        "    │           │                                                  │",
        "    │    ┌──────┴──────┐                                           │",
        "    │    │   INVOICE   │  ← What was billed?                       │",
        "    │    │ (Aggregated)│    (invoice_total = Σ line_totals)         │",
        "    │    └──────┬──────┘                                           │",
        "    │           │                                                  │",
        "    │           │ Validate vendor_id + currency                    │",
        "    │           │                                                  │",
        "    │    ┌──────┴──────┐                                           │",
        "    │    │  GOODS REC  │  ← What was received?                     │",
        "    │    │  NOTE (GRN) │    (grn_number, grn_date confirm receipt) │",
        "    │    └─────────────┘                                           │",
        "    │                                                              │",
        "    │    Decision: Pay only if ORDER = BILLED = RECEIVED           │",
        "    └──────────────────────────────────────────────────────────────┘",
    ], 6, "Industry standard: verify PO ↔ Invoice ↔ GRN before payment release")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 7: Classification Engine — Decision Tree
    # ═══════════════════════════════════════════════════════════════════
    add_diagram_slide(prs, "Classification Engine — Decision Flow", [
        "                    ┌──────────────────┐",
        "                    │ Matched Invoice  │",
        "                    │   + PO Record    │",
        "                    └────────┬─────────┘",
        "                             │",
        "                    ┌────────▼─────────┐",
        "                    │ PO found?        │──── NO ──▶ MISSING_PO (conf: 1.0)",
        "                    └────────┬─────────┘",
        "                             │ YES",
        "                    ┌────────▼─────────┐",
        "                    │ diff = inv - po  │",
        "                    │ |diff| > tol?    │──── NO ──▶ MATCHED (no mismatch)",
        "                    └────────┬─────────┘",
        "                             │ YES",
        "                    ┌────────▼─────────┐",
        "                    │ ratio = |diff|   │",
        "                    │  / min(inv, po)  │",
        "                    └────────┬─────────┘",
        "                             │",
        "              ┌──────────────┼──────────────┐",
        "              │              │              │",
        "       5% ≤ ratio ≤ 25%   ratio > 20%   else (< 5%)",
        "              │              │              │",
        "              ▼              ▼              ▼",
        "        TAX_MISCODE   QUANTITY_VAR    PRICE_VARIANCE",
        "        (conf: 0.80)  (conf: 0.75)   (conf: 0.70)",
    ], 7, "Priority-ordered rules prevent multi-classification; confidence enables triage")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 8: LangGraph Agentic Workflow — State Machine
    # ═══════════════════════════════════════════════════════════════════
    add_diagram_slide(prs, "LangGraph Agentic Workflow — State Machine", [
        "  ┌─────────────────────────────────────────────────────────────────────┐",
        "  │                    ReconciliationState (TypedDict)                    │",
        "  │   invoices | plan | matches | mismatches | emails | approved | step  │",
        "  └─────────────────────────────────────────────────────────────────────┘",
        "",
        "        ┌──────────┐      ┌──────────┐      ┌──────────────┐",
        "  START─▶│ PLANNER  │─────▶│ MATCHER  │─────▶│   DISPUTE    │",
        "        │          │      │ (tool)   │      │  GENERATOR   │",
        "        │ Analyze  │      │ Match +  │      │  LLM/Template│",
        "        │ & plan   │      │ Classify │      │  email draft │",
        "        └──────────┘      └──────────┘      └──────┬───────┘",
        "                                                     │",
        "                                            ┌────────▼────────┐",
        "                                            │  HUMAN APPROVAL │",
        "                                            │     GATE        │",
        "                                            │  ┌───┐  ┌───┐  │",
        "                                            │  │ Y │  │ N │  │",
        "                                            │  └─┬─┘  └─┬─┘  │",
        "                                            └────┼──────┼────┘",
        "                                                 │      │",
        "                                            ┌────▼──┐ ┌─▼──────┐",
        "                                            │  END  │ │DISCARD │",
        "                                            │(send) │ │+ audit │",
        "                                            └───────┘ └────────┘",
    ], 8, "LangGraph: graph-structured orchestration with explicit state, conditional edges, HITL interrupt")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 9: Human-in-the-Loop Governance Framework
    # ═══════════════════════════════════════════════════════════════════
    add_content_slide(prs, "Human-in-the-Loop (HITL) Governance", [
        "## Why HITL? (Industry Context)",
        "  • HITL reduces hallucination-related errors by 96% (AnyReach, 2025)",
        "  • EU AI Act Article 14: Human oversight mandatory for high-risk AI",
        "  • Finance domain: No automated external action without approval",
        "",
        "## SmartPay AP HITL Implementation",
        "  • Approval Gate: Every email requires explicit y/n before send",
        "  • Confidence Display: Each draft shows classification confidence %",
        "  • Rejection Audit: Discarded drafts logged with reason for compliance",
        "  • Progressive Autonomy: Design supports future confidence-based auto-routing",
        "",
        "## Guardrail Stack",
        "  • Layer 1: Input Validation — regex enforcement (^INV\\d+$)",
        "  • Layer 2: Tool Whitelist — only [match_invoices, classify, generate_email]",
        "  • Layer 3: State Enforcement — prevent_auto_send() checks workflow step",
        "  • Layer 4: GuardrailViolation exception halts + logs any breach",
    ], 9)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 10: Evaluation & Model Metrics
    # ═══════════════════════════════════════════════════════════════════
    add_table_slide(prs, "Evaluation Framework & Metrics", 
        ["Metric", "Method", "Purpose"],
        [
            ["Precision (per-class)", "sklearn classification_report", "False positive control"],
            ["Recall (per-class)", "sklearn classification_report", "Coverage completeness"],
            ["F1 (per-class)", "Harmonic mean of P & R", "Balanced accuracy"],
            ["Macro F1", "Unweighted avg across 4 classes", "Overall system quality"],
            ["Support", "Count per class in ground truth", "Class imbalance detection"],
            ["Alignment", "Inner-join on invoice_id", "Handle pred/truth size mismatch"],
        ], 10)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 11: Technology Architecture Decisions (ADRs)
    # ═══════════════════════════════════════════════════════════════════
    add_table_slide(prs, "Architecture Decision Records (ADRs)",
        ["Decision", "Choice", "Alternative", "Rationale"],
        [
            ["Orchestration", "LangGraph", "CrewAI, AutoGen", "Explicit state, native HITL interrupts"],
            ["Classification", "Rule-based", "Random Forest, XGBoost", "82 labels insufficient; auditability"],
            ["Email Gen", "LLM + Template", "LLM-only", "Graceful degradation, no API dependency"],
            ["Matching", "ID-suffix", "Fuzzy, ML-embed", "100% precision on structured IDs"],
            ["Data Processing", "pandas", "PySpark, Polars", "Right-sized for <1K records"],
            ["State Mgmt", "TypedDict", "Pydantic, dataclass", "LangGraph native, zero overhead"],
        ], 11)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 12: Production Deployment Topology
    # ═══════════════════════════════════════════════════════════════════
    add_diagram_slide(prs, "Production Deployment — Multi-Cloud Topology", [
        "  ┌────────────────── AZURE (Primary Region) ────────────────────────────┐",
        "  │                                                                       │",
        "  │  ┌─────────────┐    ┌──────────────────┐    ┌───────────────────┐   │",
        "  │  │ Blob Storage│───▶│ Azure Functions  │───▶│  Azure OpenAI     │   │",
        "  │  │ (CSV ingest)│    │ (Matching Engine)│    │  (Email drafting) │   │",
        "  │  └─────────────┘    └────────┬─────────┘    └───────────────────┘   │",
        "  │                              │                                        │",
        "  │                     ┌────────▼─────────┐                             │",
        "  │                     │  Cosmos DB        │                             │",
        "  │                     │  (Workflow State) │                             │",
        "  │                     └──────────────────┘                              │",
        "  └───────────────────────────────────────────────────────────────────────┘",
        "                              │ Replication",
        "  ┌───────────────────────────▼───────────────────────────────────────────┐",
        "  │                    AWS (DR Region)                                      │",
        "  │  S3 (replicated data) → Lambda (backup) → Bedrock (LLM fallback)      │",
        "  │  DynamoDB (state replica)                                               │",
        "  └───────────────────────────────────────────────────────────────────────┘",
        "",
        "  Shared: Azure AD/IAM | Key Vault/Secrets Manager | CI/CD | Monitor/CW",
        "  Estimated cost: $50–$100/mo (serverless, scales to zero when idle)",
    ], 12, "Active-passive multi-cloud: Azure primary (enterprise agreement), AWS for DR + Bedrock fallback")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 13: Observability, SRE & Compliance
    # ═══════════════════════════════════════════════════════════════════
    add_content_slide(prs, "Observability, SRE & Compliance", [
        "## Observability Stack",
        "  • Structured logging (Python logging → Azure Monitor / CloudWatch)",
        "  • Metrics: match rate, mismatch volume by type, approval rate, latency",
        "  • Alerting: F1 < 0.85, rejection rate > 50%, guardrail violations",
        "",
        "## SRE Considerations",
        "  • SLO: 99.9% pipeline availability (batch, non-realtime)",
        "  • Circuit breaker on OpenAI API → template fallback within 500ms",
        "  • Idempotent processing: re-run safe (CSV inputs are immutable)",
        "",
        "## Compliance Mapping",
        "  • EU AI Act Art. 14: Human oversight ✓ (Approval Gate)",
        "  • EU AI Act Art. 11: Technical documentation ✓ (docs/architecture.md)",
        "  • SOX Section 404: Internal controls ✓ (3-way match + audit log)",
        "  • GDPR: No personal data in pipeline (vendor = business entity)",
        "  • AI Literacy (Art. 4): README + architecture docs for operators",
    ], 13)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 14: Scalability Path & Future Architecture
    # ═══════════════════════════════════════════════════════════════════
    add_table_slide(prs, "Scalability Path & Future Architecture",
        ["Dimension", "Current (MVP)", "Phase 2 (6mo)", "Phase 3 (12mo)"],
        [
            ["Volume", "300 invoices", "10K/day (Polars)", "100K/day (Spark)"],
            ["Matching", "ID-suffix (deterministic)", "Fuzzy + embedding", "ML ensemble"],
            ["Classification", "Rule-based (4 types)", "ML + rules hybrid", "Deep learning"],
            ["Ingestion", "CSV upload", "Event-driven (Blob trigger)", "Real-time stream"],
            ["Approval", "Console y/n", "Web UI + Slack", "Confidence auto-route"],
            ["LLM", "GPT-3.5 + template", "GPT-4o + fine-tuned", "Custom fine-tune"],
            ["Observability", "Structured logs", "Grafana dashboards", "Anomaly detection"],
        ], 14)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 15: Summary & Architecture Principles
    # ═══════════════════════════════════════════════════════════════════
    add_content_slide(prs, "Summary — Architecture Principles", [
        "## Delivered Capabilities",
        "  ✓ End-to-end automated matching (<1s for 300 invoices)",
        "  ✓ 4-class mismatch taxonomy with confidence scoring",
        "  ✓ LangGraph agentic orchestration with typed state",
        "  ✓ Dual-mode email generation (LLM + template fallback)",
        "  ✓ Multi-layer guardrails + human approval gate",
        "  ✓ 96 unit tests covering all modules",
        "",
        "## Guiding Principles",
        "  1. Determinism First — Rules before ML (auditable, reproducible)",
        "  2. Human-in-the-Loop — AI proposes, humans approve (EU AI Act ready)",
        "  3. Graceful Degradation — Template fallback, error isolation",
        "  4. Right-Sized Complexity — No over-engineering for current scale",
        "  5. Production-Ready Design — Guardrails, observability, DR from day one",
        "  6. Progressive Autonomy — Architecture supports reducing human load over time",
        "",
        "## Next Step",
        "  Deploy MVP → Collect labelled data → Train ML classifier → Expand to full AP",
    ], 15)

    # ─── Save ─────────────────────────────────────────────────────────
    output_path = r"d:\Case Study\SmartPayAP\docs\SmartPay_AP_Architecture_Advanced.pptx"
    prs.save(output_path)
    print(f"\n✓ Advanced presentation saved to: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
