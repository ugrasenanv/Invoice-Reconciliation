"""Generate SmartPay AP presentation as .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Constants
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xAB)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
GREEN = RGBColor(0x27, 0xAE, 0x60)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BLUE

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    p2.alignment = PP_ALIGN.CENTER
    return slide


def add_content_slide(prs, title, bullets, subnotes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_top = Pt(12)
    tf.margin_left = Pt(20)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)
        if bullet.startswith("  "):
            p.level = 1
            p.font.size = Pt(14)

    if subnotes:
        p = tf.add_paragraph()
        p.space_before = Pt(16)
        p.text = subnotes
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    return slide


def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_top = Pt(12)
    tf.margin_left = Pt(20)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    left = Inches(0.5)
    top = Inches(1.4)
    width = Inches(9.0)
    height = Inches(0.4) * num_rows

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths evenly
    col_width = int(width / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_width

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_TEXT
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "SmartPay AP",
        "AI-Powered Invoice Reconciliation for Acme Manufacturing\n\nAI Architect Case Study"
    )

    # Slide 2: Problem Statement
    add_content_slide(prs, "Problem Statement", [
        "Acme Manufacturing's AP team manually reconciles 300+ invoices per cycle",
        "",
        "• Manual matching takes 8–12 hours per reconciliation cycle",
        "• Missed mismatches cause revenue leakage from undetected variances",
        "• Inconsistent follow-up leads to delayed dispute communications",
        "• No audit trail creates compliance gaps",
        "",
        "Goal: Reduce reconciliation time by 80%+ while maintaining",
        "human oversight on all external actions"
    ])

    # Slide 3: Solution Overview
    add_content_slide(prs, "Solution Overview", [
        "Three-Layer Architecture:",
        "",
        "1. DATA PIPELINE — Ingest, validate, aggregate invoice/PO data",
        "2. MATCHING & CLASSIFICATION — Rule-based mismatch detection (4 types)",
        "3. AGENTIC ORCHESTRATION — LangGraph workflow with Human-in-the-Loop",
        "",
        "End-to-End Flow:",
        "CSV Data → Matching Model → Classification → Agent Workflow → Human Approval",
        "",
        "Key Principle: The AI agent plans and proposes; humans approve and act."
    ])

    # Slide 4: Architecture Diagram
    add_content_slide(prs, "High-Level Architecture", [
        "┌─────────────────────────────────────────────────────────┐",
        "│  DATA LAYER    →    CLASSIFICATION    →    AGENTIC LAYER │",
        "└─────────────────────────────────────────────────────────┘",
        "",
        "• Data Layer: CSV loading, validation, date parsing, aggregation",
        "• Classification: Matcher (ID-suffix join) + Classifier (4 rules)",
        "• Agent: LangGraph StateGraph with 4 nodes + conditional edges",
        "• Guardrails: Input validation, tool whitelist, HITL gate",
        "",
        "All layers connected through clean Python interfaces (no tight coupling)"
    ])

    # Slide 5: Data Pipeline
    add_content_slide(prs, "Data Pipeline", [
        "Input: AcmeMini Dataset (3 CSV files)",
        "  • invoices.csv — 1,500 line items across 300 invoices",
        "  • po_grn.csv — 300 purchase orders with GRN data",
        "  • labelled_mismatches.csv — 82 ground-truth labels",
        "",
        "Processing Steps:",
        "  1. Load & validate columns (schema enforcement)",
        "  2. Parse dates (DD-MM-YYYY → datetime)",
        "  3. Handle nulls/invalid numerics (log warning & exclude)",
        "  4. Aggregate line items → one total per invoice_id",
        "",
        "Design: Fail fast on schema issues, graceful on data quality"
    ])

    # Slide 6: Matching Strategy
    add_content_slide(prs, "Matching Strategy", [
        "Approach: Numeric ID Suffix Extraction",
        "",
        "  INV0001 → extract '0001' → match to PO0001",
        "  INV0042 → extract '0042' → match to PO0042",
        "",
        "Why this approach?",
        "  • AcmeMini uses clear INV000X ↔ PO000X convention",
        "  • Achieves 100% match rate on current dataset",
        "  • More reliable than fuzzy matching for structured IDs",
        "",
        "Validation Flags: vendor_match, currency_match",
        "Semantics: Left-join — all invoices appear; unmatched ones flagged"
    ])

    # Slide 7: Classification Engine (Table)
    add_table_slide(prs, "Classification Engine — Rule-Based", 
        ["Priority", "Type", "Detection Logic", "Confidence"],
        [
            ["1", "MISSING_PO", "No matching PO found", "1.0"],
            ["2", "TAX_MISCODE", "Difference is 5–25% of min total", "0.8"],
            ["3", "QUANTITY_VARIANCE", "Ratio > 20% (whole-unit delta)", "0.75"],
            ["4", "PRICE_VARIANCE", "Remaining diffs above tolerance", "0.70"],
        ]
    )

    # Slide 8: Why Rules over ML
    add_content_slide(prs, "Why Rules Over ML?", [
        "Decision: Rule-based classification over Machine Learning",
        "",
        "Rationale:",
        "  • Only 82 labelled samples — insufficient for reliable ML training",
        "  • Rules are interpretable and auditable (critical for finance)",
        "  • Thresholds are configurable without retraining",
        "  • Priority ordering prevents multi-classification",
        "",
        "Future Path:",
        "  • When >500 labelled records accumulate → train ML model",
        "  • Current rule-based metrics serve as ML benchmark",
        "  • Hybrid approach: ML predictions + rule-based confidence scoring"
    ])

    # Slide 9: Agentic Workflow
    add_content_slide(prs, "Agentic Workflow (LangGraph)", [
        "Workflow Graph:",
        "  START → Planner → Matcher → Dispute Generator → Human Approval → END",
        "",
        "Node Responsibilities:",
        "  • Planner: Analyze invoices, produce reconciliation plan",
        "  • Matcher: Invoke matching model as tool, classify mismatches",
        "  • Dispute Generator: Draft emails per mismatch (LLM or template)",
        "  • Human Approval: Display draft, wait for y/n confirmation",
        "",
        "Conditional Edge:",
        "  • Approved → END (proceed with email)",
        "  • Rejected → Discard node (log rejection, drop draft)"
    ])

    # Slide 10: Email Generation
    add_content_slide(prs, "Dispute Email Generation", [
        "Dual-Mode Strategy:",
        "",
        "  • LLM Mode (GPT-3.5-turbo): Natural, contextual language",
        "  • Template Mode (fallback): Consistent, reliable, no API needed",
        "",
        "Email Contents:",
        "  • Professional greeting addressed to vendor",
        "  • Invoice ID + PO number references",
        "  • Mismatch type and exact discrepancy amount",
        "  • Resolution request per type:",
        "    — PRICE/QUANTITY → credit note or corrected invoice",
        "    — TAX_MISCODE → tax calculation correction",
        "    — MISSING_PO → provide PO reference"
    ])

    # Slide 11: Guardrails & Safety (Table)
    add_table_slide(prs, "Guardrails & Responsible AI",
        ["Guardrail", "What It Prevents"],
        [
            ["Input Validation", "Malformed invoice IDs (must match INV\\d+)"],
            ["Tool Whitelist", "Agent calling unregistered/arbitrary tools"],
            ["Auto-Send Prevention", "Emails sent without human review"],
            ["Rejection Logging", "Full audit trail for discarded drafts"],
        ]
    )

    # Slide 12: Technology Stack (Table)
    add_table_slide(prs, "Technology Stack",
        ["Technology", "Purpose", "Why Chosen"],
        [
            ["Python 3.10+", "Primary language", "Rich ML/data ecosystem"],
            ["pandas", "Data processing", "Ideal for <1K row datasets"],
            ["scikit-learn", "Evaluation metrics", "Industry standard P/R/F1"],
            ["LangGraph", "Agent orchestration", "Explicit HITL, state control"],
            ["OpenAI GPT-3.5", "Email generation", "Cost-effective, quality output"],
            ["pytest", "Testing (96 tests)", "Concise, parametrize support"],
        ]
    )

    # Slide 13: Production Deployment
    add_content_slide(prs, "Production Deployment — Multi-Cloud", [
        "Azure (Primary):",
        "  • Blob Storage → Azure Functions → Azure OpenAI",
        "  • Cosmos DB for workflow state",
        "  • Azure AD for identity & access",
        "",
        "AWS (Secondary / DR):",
        "  • S3 → Lambda → Amazon Bedrock (LLM fallback)",
        "  • DynamoDB for state replication",
        "",
        "Shared Services:",
        "  • Key Vault / Secrets Manager for API keys",
        "  • CI/CD via GitHub Actions / Azure DevOps",
        "",
        "Estimated Cost: ~$50–$100/month (serverless, scales to zero)"
    ])

    # Slide 14: Scalability & Roadmap
    add_table_slide(prs, "Scalability & Future Roadmap",
        ["Dimension", "Current (AcmeMini)", "Production Future"],
        [
            ["Data Volume", "300 invoices", "10K+/day → Spark/Polars"],
            ["Matching", "ID-suffix pairing", "Fuzzy + ML matching"],
            ["Classification", "Rules (82 labels)", "ML model (500+ labels)"],
            ["Processing", "Batch (<1 sec)", "Event-driven (real-time)"],
            ["Currencies", "Single", "Multi-currency + FX rates"],
            ["Approval UI", "Console y/n", "Web UI + approval chains"],
        ]
    )

    # Slide 15: Summary
    add_content_slide(prs, "Summary & Key Takeaways", [
        "What SmartPay AP Delivers:",
        "",
        "  ✓ Automated matching — 300 invoices matched in <1 second",
        "  ✓ Intelligent classification — 4 mismatch types, configurable rules",
        "  ✓ Agentic workflow — End-to-end orchestration with LangGraph",
        "  ✓ Human oversight — No external action without explicit approval",
        "  ✓ Production-ready — Guardrails, audit logging, template fallback",
        "  ✓ Extensible — Clear path from rules to ML as data grows",
        "",
        "Design Principles:",
        "  • Interpretability over black-box (auditable rules)",
        "  • Human-in-the-loop (AI proposes, humans approve)",
        "  • Graceful degradation (template fallback, error tolerance)",
        "  • Right-sized technology (no over-engineering)"
    ])

    # Save
    output_path = r"d:\Case Study\SmartPayAP\docs\SmartPay_AP_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
